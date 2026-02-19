# filepath: /home/pablo/Documentos/Coding/AGP/Markdowner/markdowner.py
import pypandoc
import os
import argparse
import tempfile
from pdfminer.high_level import extract_text_to_fp
from io import StringIO

def convert_docx_to_markdown(docx_path, markdown_path):
    """
    Converts a DOCX file to Markdown.

    Args:
        docx_path (str): The path to the input DOCX file.
        markdown_path (str): The path to save the output Markdown file.
    """
    try:
        pypandoc.convert_file(docx_path, 'markdown', outputfile=markdown_path)
        print(f"Successfully converted '{docx_path}' to '{markdown_path}'")
    except Exception as e:
        print(f"Error converting DOCX to Markdown: {e}")

def convert_pptx_to_markdown(pptx_path, markdown_path):
    """
    Converts a PPTX file to Markdown by extracting text content.
    
    Args:
        pptx_path (str): The path to the input PPTX file.
        markdown_path (str): The path to save the output Markdown file.
    """
    try:
        from pptx import Presentation
        
        # Load the presentation
        prs = Presentation(pptx_path)
        
        # Extract text from all slides
        markdown_content = []
        
        for i, slide in enumerate(prs.slides, 1):
            markdown_content.append(f"# Slide {i}\n")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    # Add text with proper markdown formatting
                    text = shape.text.strip()
                    if text:
                        # Simple heuristic: if text is short and doesn't end with punctuation, treat as heading
                        if len(text) < 100 and not text.endswith(('.', '!', '?', ':')):
                            markdown_content.append(f"## {text}\n")
                        else:
                            markdown_content.append(f"{text}\n")
            
            markdown_content.append("\n---\n")  # Slide separator
        
        # Write to markdown file
        with open(markdown_path, 'w', encoding='utf-8') as f:
            f.write('\n'.join(markdown_content))
        
        print(f"Successfully converted '{pptx_path}' to '{markdown_path}'")
        
    except ImportError:
        print("Error: python-pptx library not installed. Install it with: pip install python-pptx")
    except Exception as e:
        print(f"Error converting PPTX to Markdown: {e}")

def convert_pdf_to_markdown(pdf_path, markdown_path, force_ocr=False):
    """
    Converts a PDF file to Markdown-friendly text. 
    Prioritizes 'pdftotext' with layout preservation. 
    If result is too short (likely a scan) or force_ocr is True, falls back to OCR.
    """
    import subprocess
    import shutil

    extracted_success = False
    
    # Method 1: Try pdftotext (System CLI) - Best for preserving visual layout (tables, columns)
    if not force_ocr and shutil.which('pdftotext'):
        try:
            subprocess.run(['pdftotext', '-layout', pdf_path, markdown_path], check=True)
            
            # Check if the extracted file has significant content
            if os.path.exists(markdown_path) and os.path.getsize(markdown_path) > 100:
                print(f"Successfully converted '{pdf_path}' to '{markdown_path}' using pdftotext")
                extracted_success = True
            else:
                print("Warning: pdftotext yielded very little text. PDF might be a scan.")
        except subprocess.CalledProcessError as e:
            print(f"Warning: pdftotext failed ({e}).")

    # Method 2: Fallback to OCR if pdftotext failed or for manual override
    if not extracted_success:
        try:
            print(f"Attempting OCR extraction for '{pdf_path}'...")
            from pdf2image import convert_from_path
            import pytesseract
            from PIL import Image

            # We use a temp directory to avoid cluttering with page images
            with tempfile.TemporaryDirectory() as temp_dir:
                pages = convert_from_path(pdf_path)
                full_text = []
                for i, page in enumerate(pages):
                    img_path = os.path.join(temp_dir, f'page_{i}.png')
                    page.save(img_path, 'PNG')
                    # 'por' for Portuguese, 'eng' for English. Could be parameterized.
                    page_text = pytesseract.image_to_string(Image.open(img_path), lang='por+eng')
                    full_text.append(page_text)
                
                with open(markdown_path, 'w', encoding='utf-8') as f:
                    f.write("\n\n".join(full_text))
            
            print(f"Successfully converted '{pdf_path}' to '{markdown_path}' using OCR")
            extracted_success = True

        except ImportError:
            print("Error: OCR dependencies missing. Run: sudo apt install tesseract-ocr python3-pytesseract python3-pdf2image")
        except Exception as e:
            print(f"Error during OCR: {e}")

    # Method 3: Final Fallback to pdfminer.six (Digital text only)
    if not extracted_success:
        try:
            from pdfminer.high_level import extract_text
            text = extract_text(pdf_path)
            with open(markdown_path, 'w', encoding='utf-8') as f:
                f.write(text)
            print(f"Successfully converted '{pdf_path}' to '{markdown_path}' using pdfminer fallback")
        except Exception as e:
            print(f"All extraction methods failed for '{pdf_path}': {e}")


def perform_conversion(input_file, output_file, force_ocr=False):
    """Wrapper for actual conversion logic, ensures output dir exists."""
    file_ext = os.path.splitext(input_file)[1].lower()
    if file_ext not in ['.docx', '.pdf', '.pptx']:
        print(f"Skipping '{input_file}': Unsupported file type '{file_ext}'.")
        return False # Indicates not processed or error

    output_dir = os.path.dirname(output_file)
    if output_dir: # If output_file has a directory part
        os.makedirs(output_dir, exist_ok=True)

    if file_ext == '.docx':
        convert_docx_to_markdown(input_file, output_file)
    elif file_ext == '.pdf':
        print(f"Converting '{input_file}' to '{output_file}'...")
        convert_pdf_to_markdown(input_file, output_file, force_ocr=force_ocr)
    elif file_ext == '.pptx':
        convert_pptx_to_markdown(input_file, output_file)
    return True # Indicates processed


def main():
    parser = argparse.ArgumentParser(description="Convert DOCX, PDF, and PPTX files to Markdown.")
    parser.add_argument("input_path", help="Path to the input DOCX/PDF/PPTX file or a directory.")
    parser.add_argument("-o", "--output", help="Path to the output Markdown file or directory.")
    parser.add_argument("--ocr", action="store_true", help="Force OCR extraction for PDFs (slower, but works for scans).")
    args = parser.parse_args()
    input_path = args.input_path
    output_option = args.output
    force_ocr = args.ocr

    if not os.path.exists(input_path):
        print(f"Error: Input path '{input_path}' not found.")
        return

    files_to_convert = [] # List of (input_file, output_file) tuples

    if os.path.isfile(input_path):
        file_path_no_ext, input_ext = os.path.splitext(input_path)
        if input_ext.lower() not in ['.docx', '.pdf', '.pptx']:
            print(f"Error: Input file '{input_path}' is not a .docx, .pdf, or .pptx file.")
            return
            
        output_file = ""
        if not output_option:
            output_file = file_path_no_ext + ".md"
        elif os.path.isdir(output_option): # -o is an existing directory
            output_file = os.path.join(output_option, os.path.basename(file_path_no_ext) + ".md")
        else: # -o is a specific file path
            output_file = output_option
        files_to_convert.append((input_path, output_file))

        
    elif os.path.isdir(input_path):
        if output_option and os.path.exists(output_option) and os.path.isfile(output_option):
            print(f"Error: If input is a directory, output path '{output_option}' must be a directory, not a file.")
            return
        
        found_supported_file = False
        for item_name in os.listdir(input_path):
            item_full_path = os.path.join(input_path, item_name)
            if os.path.isfile(item_full_path):
                item_path_no_ext, item_ext = os.path.splitext(item_full_path)
                if item_ext.lower() not in ['.docx', '.pdf', '.pptx']:
                    continue
                found_supported_file = True
                
                output_file = ""
                if not output_option: # Save next to original
                    output_file = item_path_no_ext + ".md"
                else: # output_option is a directory (either existing or to be created by perform_conversion)
                    output_file = os.path.join(output_option, os.path.basename(item_path_no_ext) + ".md")
                files_to_convert.append((item_full_path, output_file))
        
        if not found_supported_file:
            print(f"No .docx, .pdf, or .pptx files found in directory '{input_path}'.")
            return
    else:
        print(f"Error: Input path '{input_path}' is not a valid file or directory.")
        return

    if not files_to_convert:
        print("No files selected for conversion.")
        return

    for input_f, output_f in files_to_convert:
        perform_conversion(input_f, output_f, force_ocr=force_ocr)

if __name__ == "__main__":
    main()
